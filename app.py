import io
import json
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

app = Flask(__name__)
CORS(app)


def delete_slide(prs, index):
    """Διαγράφει ασφαλώς μια διαφάνεια από το Presentation βάσει του index της."""
    slide_id = prs.slides[index].slide_id
    for element in prs.slides._sldIdLst:
        if element.id == slide_id:
            rId = element.rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(element)
            break


def get_best_content_layout(prs):
    """🎯 ΠΛΗΡΩΣ ΔΥΝΑΜΙΚΗ ΕΥΡΕΣΗ:

    Ψάχνει στα Master Layouts του PowerPoint για το standard layout
    που περιέχει ΚΑΙ Τίτλο ΚΑΙ Body (Bullets), χωρίς να βασίζεται στη 3η διαφάνεια.
    """
    # 1. Σάρωση όλων των διαθέσιμων Slide Layouts του PPTX template
    for layout in prs.slide_layouts:
        has_title = False
        has_body = False
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                has_title = True
            elif ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                has_body = True

        # Αν το layout έχει και τίτλο και bullets, αυτό είναι το Main Content Layout!
        if has_title and has_body:
            return layout

    # 2. Fallback αν δεν βρεθεί μέσω placeholders:
    # Σαρώνει τις υπάρχουσες διαφάνειες εκτός της 1ης (Cover)
    if len(prs.slides) > 1:
        for slide in prs.slides[1:]:
            # Αν η διαφάνεια δεν είναι σκοτεινή/Agenda (έλεγχος αν έχει body)
            for shape in slide.shapes:
                if (
                    shape.is_placeholder
                    and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
                ):
                    return slide.slide_layout

    # 3. Default Fallback
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]


def process_cover_slide(slide, deck_title, deck_subtitle=""):
    """🎯 ΕΞΩΦΥΛΛΟ: Βάζει ΜΟΝΟ τον τίτλο (και υπότιτλο) και καθαρίζει ΟΛΑ τα άλλα κείμενα."""
    title_done = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type

            # Τίτλος
            if (
                ph_type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE]
                and not title_done
            ):
                if len(tf.paragraphs) > 0:
                    p = tf.paragraphs[0]
                    if len(p.runs) > 0:
                        p.runs[0].text = deck_title
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = deck_title
                title_done = True

            # Υπότιτλος
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                if deck_subtitle and len(tf.paragraphs) > 0:
                    p = tf.paragraphs[0]
                    if len(p.runs) > 0:
                        p.runs[0].text = deck_subtitle
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = deck_subtitle
                else:
                    tf.text = ""

            # Σβήνουμε οτιδήποτε άλλο (Body/Bullets) στο Cover Slide
            elif ph_type not in [
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.HEADER,
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.DATE,
            ]:
                tf.text = ""

        else:
            if not title_done and shape.top < 4000000:
                tf.text = deck_title
                title_done = True
            else:
                tf.text = ""


def process_slide_text(slide, title_text, bullets):
    """Εντοπίζει αυστηρά μόνο τον Τίτλο και το Κύριο Σώμα (Body) στις διαφάνειες περιεχομένου."""
    title_shape = None
    body_shape = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type

            if ph_type in [
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.HEADER,
                PP_PLACEHOLDER.DATE,
            ]:
                continue

            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                title_shape = shape
            elif ph_type in [
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.OBJECT,
            ]:
                body_shape = shape

    if not title_shape or not body_shape:
        valid_text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if shape.top > 5943600:  # Αγνοούμε footers
                continue

            valid_text_shapes.append(shape)

        valid_text_shapes.sort(key=lambda s: s.top)

        if valid_text_shapes:
            if not title_shape:
                title_shape = valid_text_shapes[0]
            if not body_shape and len(valid_text_shapes) > 1:
                body_shape = valid_text_shapes[1]

    # --- ΕΦΑΡΜΟΓΗ ΚΕΙΜΕΝΟΥ ΣΤΟΝ ΤΙΤΛΟ ---
    if title_shape and title_shape.has_text_frame:
        tf = title_shape.text_frame
        if title_text and len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
            if len(p.runs) > 0:
                p.runs[0].text = title_text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = title_text

    # --- ΕΦΑΡΜΟΓΗ ΚΕΙΜΕΝΟΥ ΣΤΑ BULLETS ---
    if body_shape and body_shape.has_text_frame and bullets:
        tf = body_shape.text_frame

        for p_idx, paragraph in enumerate(tf.paragraphs):
            if p_idx < len(bullets):
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = bullets[p_idx]
                    for r in paragraph.runs[1:]:
                        r.text = ""
                else:
                    paragraph.text = bullets[p_idx]
            else:
                paragraph.text = ""

        if len(bullets) > len(tf.paragraphs):
            for b_idx in range(len(tf.paragraphs), len(bullets)):
                p = tf.add_paragraph()
                p.text = bullets[b_idx]


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"}), 200


@app.route("/inject", methods=["POST"])
def inject_text():
    try:
        if "template" not in request.files or "data" not in request.form:
            return jsonify({"error": "Missing template file or data JSON"}), 400

        template_file = request.files["template"]
        data_json = json.loads(request.form["data"])

        deck_title = data_json.get("title", "")
        deck_subtitle = data_json.get("subtitle", "")
        slides_data = data_json.get("slides", [])

        prs = Presentation(template_file)

        # 1. ΧΕΙΡΙΣΜΟΣ COVER SLIDE (Διαφάνεια 1 - Index 0)
        if len(prs.slides) > 0:
            process_cover_slide(prs.slides[0], deck_title, deck_subtitle)

        # 2. ΔΥΝΑΜΙΚΗ ΕΥΡΕΣΗ LAYOUT ΠΕΡΙΕΧΟΜΕΝΟΥ (Από τα Master Layouts)
        content_layout = get_best_content_layout(prs)

        # 3. ΚΑΘΑΡΙΖΟΥΜΕ/ΔΙΑΓΡΑΦΟΥΜΕ ΟΛΕΣ ΤΙΣ ΠΑΛΙΕΣ ΔΙΑΦΑΝΕΙΕΣ ΜΕΤΑ ΤΟ COVER
        while len(prs.slides) > 1:
            delete_slide(prs, 1)

        # 4. ΔΗΜΙΟΥΡΓΟΥΜΕ ΝΕΕΣ ΔΙΑΦΑΝΕΙΕΣ ΓΙΑ ΚΑΘΕ SLIDE ΤΟΥ GEMINI
        for slide_info in slides_data:
            slide = prs.slides.add_slide(content_layout)
            process_slide_text(
                slide,
                slide_info.get("title", ""),
                slide_info.get("bullets", []),
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return send_file(
            output,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name="presentation_updated.pptx",
        )

    except Exception as e:
        print(f"Error processing PPTX: {e!s}")
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
